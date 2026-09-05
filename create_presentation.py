"""
ForecastBI - Professional PowerPoint Generator
Clean, tight, executive-grade 8-slide deck.
"""

import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ─── PATHS ─────────────────────────────────────────────────────────────────
OUT = Path(r"d:\A- Commodity Price Coordinator\antigravity projects\ForecastBI\ForecastBI_Presentation.pptx")
REPORTS = Path(r"d:\A- Commodity Price Coordinator\antigravity projects\ForecastBI\reports\ForecastBI_Presentation.pptx")
ARTIFACT = Path(r"C:\Users\sa\.gemini\antigravity-ide\brain\a572f10e-0c4b-4b65-a282-b8dfa98da6da\ForecastBI_Presentation.pptx")

# ─── PALETTE ────────────────────────────────────────────────────────────────
NAVY      = RGBColor(11,  23, 44)    # slide background / title bg
SLATE     = RGBColor(248,250,252)    # body slide bg
WHITE     = RGBColor(255,255,255)
BLUE      = RGBColor(37,  99, 235)   # primary accent
SKY       = RGBColor(56, 182, 255)   # highlight blue
TEAL      = RGBColor(5,  150,105)    # green accent
AMBER     = RGBColor(217,119,  6)    # amber
INK       = RGBColor(15,  23, 42)    # dark text
MUTED     = RGBColor(71,  85,105)    # secondary text
RULE      = RGBColor(203,213,225)    # divider lines
CARD_BG   = RGBColor(255,255,255)
CARD_BDR  = RGBColor(226,232,240)

# ─── HELPERS ────────────────────────────────────────────────────────────────
def px(v): return Inches(v)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, line_color=None, line_pt=0, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, px(l), px(t), px(w), px(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_pt)
    else:
        sh.line.fill.background()
    return sh

def textbox(slide, text, l, t, w, h,
            size=11, bold=False, color=INK, align=PP_ALIGN.LEFT,
            italic=False, wrap=True, font="Segoe UI", spacing_before=0):
    tb = slide.shapes.add_textbox(px(l), px(t), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(spacing_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb

def multiline(slide, lines, l, t, w, h, base_size=11, default_color=INK, wrap=True):
    """
    lines = list of (text, size, bold, color, italic, spacing_before, align, font)
    """
    tb = slide.shapes.add_textbox(px(l), px(t), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    first = True
    for item in lines:
        text, size, bold, color, italic, sp_before, align, font = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(sp_before)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb

def hline(slide, l, t, w, color=RULE):
    """Thin horizontal rule."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(l), px(t), px(w), Pt(1.2))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()

def section_label(slide, text, l=0.55, t=0.37):
    textbox(slide, text.upper(), l, t, 12, 0.25,
            size=8.5, bold=True, color=BLUE, font="Segoe UI")

def slide_title(slide, title, subtitle=None, t_top=0.55, light=False):
    title_color = WHITE if light else INK
    sub_color   = RGBColor(147,197,253) if light else MUTED
    textbox(slide, title, 0.55, t_top, 12.2, 0.52,
            size=22, bold=True, color=title_color, font="Segoe UI")
    if subtitle:
        textbox(slide, subtitle, 0.55, t_top + 0.5, 12.2, 0.35,
                size=11, bold=False, color=sub_color, font="Segoe UI")

def pill(slide, text, l, t, w, h, bg_color, text_color=WHITE, size=10):
    rect(slide, l, t, w, h, bg_color, rounded=True)
    textbox(slide, text, l + 0.08, t + 0.04, w - 0.16, h - 0.05,
            size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER, font="Segoe UI")

# ─── CARD FACTORY ───────────────────────────────────────────────────────────
def card(slide, l, t, w, h, header, body_lines, accent=BLUE, bg=CARD_BG, bdr=CARD_BDR):
    """
    Draws a rounded card with a top accent strip, a header, and body lines.
    body_lines: list of strings (auto-bullet prefix optional).
    """
    # Shadow layer
    rect(slide, l+0.03, t+0.04, w, h, RGBColor(220,228,240), rounded=True)
    # Card surface
    rect(slide, l, t, w, h, bg, bdr, 0.8, rounded=True)
    # Accent strip
    rect(slide, l, t, w, 0.055, accent, rounded=False)

    # Header
    textbox(slide, header, l+0.17, t+0.12, w-0.22, 0.32,
            size=12.5, bold=True, color=INK, font="Segoe UI")

    # Body
    body_text = "\n".join(body_lines)
    textbox(slide, body_text, l+0.17, t+0.46, w-0.22, h-0.52,
            size=10.5, bold=False, color=MUTED, font="Segoe UI", wrap=True)

def step_row(slide, num, title, desc, t, accent=BLUE):
    """Horizontal row with numbered step, title, description. Full width."""
    W = 12.23
    rect(slide, 0.55, t, W, 0.82, CARD_BG, CARD_BDR, 0.8, rounded=True)
    # Number circle
    rect(slide, 0.68, t+0.13, 0.50, 0.50, accent, rounded=True)
    textbox(slide, num, 0.68, t+0.13, 0.50, 0.50,
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Segoe UI")
    # Title
    textbox(slide, title, 1.30, t+0.10, 2.40, 0.30,
            size=12, bold=True, color=INK, font="Segoe UI")
    # Desc
    textbox(slide, desc, 1.30, t+0.40, 11.00, 0.30,
            size=10.5, bold=False, color=MUTED, font="Segoe UI")

# ─── BUILD PRESENTATION ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  ·  Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)

# Full-height left accent panel
rect(s, 0, 0, 0.30, 7.5, BLUE)

# Subtle vertical gradient bar
rect(s, 0.30, 0, 0.025, 7.5, RGBColor(30,60,120))

multiline(s, [
    ("ForecastBI", 40, True,  WHITE,                  False, 0,  PP_ALIGN.LEFT, "Segoe UI"),
    ("Pakistan Commodity Price Intelligence & Forecasting", 18, False, SKY, False, 10, PP_ALIGN.LEFT, "Segoe UI"),
    ("", 8, False, WHITE, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("End-to-End ML Pipeline  ·  Prophet vs. XGBoost  ·  Power BI Analytics", 11, False,
     RGBColor(148,163,184), False, 0, PP_ALIGN.LEFT, "Segoe UI"),
], 0.8, 2.2, 11.8, 3.5)

# Bottom stats strip
rect(s, 0.30, 6.55, 13.03, 0.95, RGBColor(18,32,60))
stats = [
    ("5 Years of Data", 1.0),
    ("7 Commodities", 4.0),
    ("17 Cities",      7.0),
    ("2 ML Models",   10.0),
]
for label, x in stats:
    textbox(s, label, x, 6.65, 2.2, 0.45,
            size=11.5, bold=True, color=SKY, align=PP_ALIGN.CENTER, font="Segoe UI")

# Thin divider above stats
rect(s, 0.30, 6.53, 13.03, 0.025, BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  ·  Introduction & Scope
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, SLATE)
hline(s, 0.55, 1.45, 12.23)
section_label(s, "Introduction")
slide_title(s, "What Is ForecastBI?",
            "An automated price-intelligence platform built on Pakistan Bureau of Statistics (PBS) data.", t_top=0.55)

# 3 cards across bottom
card(s, 0.55, 1.60, 3.90, 5.60,
     "The Context",
     ["Food price volatility directly impacts Pakistan's households, especially for perishables like Tomatoes, Onions, and Potatoes.",
      "",
      "Pulses (Moong, Masoor, Gram, Mash) are core dietary protein sources. Any prolonged spike triggers nutritional insecurity."],
     BLUE)

card(s, 4.67, 1.60, 3.90, 5.60,
     "The Objective",
     ["Convert 5+ years of scattered PBS CPI monthly bulletins into a single predictive data asset.",
      "",
      "Automatically forecast the next 12 months of retail prices per commodity across 17 commercial cities in Pakistan."],
     TEAL)

card(s, 8.78, 1.60, 4.00, 5.60,
     "Coverage",
     ["7 Key Commodities: Tomatoes, Onions, Potatoes, Pulse Moong, Masoor, Gram, Mash.",
      "",
      "17 Cities: Islamabad, Lahore, Karachi, Peshawar, Quetta, Multan, Faisalabad, Hyderabad & 9 more.",
      "",
      "41 months of validated historical records."],
     AMBER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3  ·  Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, SLATE)
hline(s, 0.55, 1.45, 12.23)
section_label(s, "Problem Statement")
slide_title(s, "Why This Problem Is Hard",
            "Three compounding challenges that make commodity price intelligence non-trivial.")

card(s, 0.55, 1.60, 3.90, 5.60,
     "Extreme Volatility",
     ["Tomatoes: single-month swing of +174% then -51%.",
      "",
      "Potatoes follow a sinusoidal harvest cycle, crashing from Rs 115 to Rs 31/kg annually.",
      "",
      "Standard regression models fail to capture these non-linear seasonal dynamics."],
     RGBColor(220, 38, 38))

card(s, 4.67, 1.60, 3.90, 5.60,
     "Fragmented Data",
     ["PBS publishes one Annexure per month across inconsistent formats: modern XLSX, legacy XLS, and read-only PDFs.",
      "",
      "No unified API. Data must be scraped and parsed from a dynamic JavaScript-rendered government web portal."],
     AMBER)

card(s, 8.78, 1.60, 4.00, 5.60,
     "Reactive Decision-Making",
     ["Procurement teams and government bodies have no forward visibility.",
      "",
      "Import quotas, reserve releases, and consumer protection measures are decided after price shocks occur — never before."],
     BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4  ·  Solution Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, SLATE)
hline(s, 0.55, 1.45, 12.23)
section_label(s, "Solution")
slide_title(s, "End-to-End Pipeline Architecture",
            "Four sequential phases from raw government bulletin to interactive dashboard.")

# Pipeline flow — 4 numbered step rows
steps = [
    ("1", "Scrape & Ingest",
     "Headless regex extraction of JS-embedded PBS data (chompjs). Dual-engine parser handles XLSX, XLS, and PDF formats.",
     BLUE),
    ("2", "Clean & Engineer Features",
     "Standardize 51 commodities × 17 cities. Build ML feature matrix: calendar transforms, autoregressive lags (1–12), rolling stats, momentum.",
     TEAL),
    ("3", "Train, Backtest & Forecast",
     "6-month holdout validation (35 train / 6 test). Benchmark Prophet vs. Recursive XGBoost. Route each commodity to its winning model. Generate 12-month forecasts.",
     BLUE),
    ("4", "Power BI Analytics Suite",
     "Export star-schema CSVs. Render 3 interactive dashboards: Executive KPIs, 12-Month Trend, Model Accuracy & City Price Intelligence.",
     AMBER),
]
for i, (num, title, desc, acc) in enumerate(steps):
    step_row(s, num, title, desc, 1.60 + i * 1.05, acc)

# Arrow connectors between rows
for i in range(3):
    rect(s, 0.875, 2.40 + i * 1.05, 0.04, 0.07, RGBColor(150,160,180))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5  ·  Data Scraping Challenges
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, SLATE)
hline(s, 0.55, 1.45, 12.23)
section_label(s, "Problems Faced")
slide_title(s, "Data Scraping Bottlenecks & How We Solved Them",
            "Real engineering challenges overcome during the data ingestion phase.")

# Layout: 2 tall cards left, 1 tall card right
card(s, 0.55, 1.60, 5.95, 2.65,
     "Problem 1 — Dynamic JS Payload (PBS Website)",
     ["Standard scrapers returned empty HTML. Download links were injected via JavaScript (const cpidata1).",
      "Solution: Extracted raw JS string with regex, parsed into Python dict with chompjs in under 2 seconds — no Selenium required."],
     RGBColor(220, 38, 38))

card(s, 0.55, 4.37, 5.95, 2.83,
     "Problem 2 — Hybrid Document Formats",
     ["Same data published as modern XLSX, legacy XLS, or read-only PDF depending on the year.",
      "Solution: Dual-engine compiler — openpyxl for Excel, pypdf + regex tokenizer for PDF line-by-line price extraction."],
     AMBER)

card(s, 6.63, 1.60, 6.10, 5.60,
     "Problem 3 — Missing Months & Data Gaps",
     ['Some months had empty URL cells ("None"), broken links, or 404 responses.',
      "",
      "Solution A — Audit Log: Script tracks a 60-month target window and writes every failure to missing_months.txt with its specific reason.",
      "",
      "Solution B — Auto-Imputation: When the official National Average cell is missing, the pipeline dynamically computes the arithmetic mean across all valid city observations for that row.",
      "",
      "Solution C — Polite Crawling: 1.5-second delay between downloads prevents IP rate-limiting."],
     TEAL)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6  ·  ML Models
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, SLATE)
hline(s, 0.55, 1.45, 12.23)
section_label(s, "Implementation — Machine Learning")
slide_title(s, "Model Design: Prophet vs. XGBoost",
            "Two complementary algorithms targeting distinct behavioral characteristics of commodity prices.")

# Left column: Prophet
rect(s, 0.55, 1.60, 5.95, 5.60, CARD_BG, CARD_BDR, 0.8, rounded=True)
rect(s, 0.55, 1.60, 5.95, 0.055, BLUE, rounded=False)

multiline(s, [
    ("Facebook Prophet", 14, True, BLUE, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("Additive/multiplicative Fourier decomposition", 10.5, False, MUTED, False, 4, PP_ALIGN.LEFT, "Segoe UI"),
    ("", 6, False, INK, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("y(t) = Trend + Seasonality + Error", 11, True, INK, True, 0, PP_ALIGN.LEFT, "Segoe UI Semibold"),
    ("", 6, False, INK, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("· Smooth Fourier yearly harmonics capture cyclical harvest patterns.", 10.5, False, MUTED, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("· Multiplicative mode for high-volatility vegetables — swings scale with price level.", 10.5, False, MUTED, False, 5, PP_ALIGN.LEFT, "Segoe UI"),
    ("· Additive mode for stable pulses — constant seasonal amplitude.", 10.5, False, MUTED, False, 5, PP_ALIGN.LEFT, "Segoe UI"),
    ("· Native upper/lower confidence intervals.", 10.5, False, MUTED, False, 5, PP_ALIGN.LEFT, "Segoe UI"),
    ("", 6, False, INK, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("Best for: Potatoes, Pulse Gram", 11, True, BLUE, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
], 0.75, 1.72, 5.70, 5.30)

# Right column: XGBoost
rect(s, 6.78, 1.60, 6.00, 5.60, CARD_BG, CARD_BDR, 0.8, rounded=True)
rect(s, 6.78, 1.60, 6.00, 0.055, TEAL, rounded=False)

multiline(s, [
    ("XGBoost — Recursive Multi-Step Regressor", 14, True, TEAL, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("Gradient-boosted decision trees with autoregressive roll", 10.5, False, MUTED, False, 4, PP_ALIGN.LEFT, "Segoe UI"),
    ("", 6, False, INK, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("Feature Matrix per forecast step:", 11, True, INK, False, 0, PP_ALIGN.LEFT, "Segoe UI Semibold"),
    ("", 6, False, INK, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("· Calendar: month, quarter, sin(2πm/12), cos(2πm/12)", 10.5, False, MUTED, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("· Lags: price at t-1, t-2, t-3, t-6, t-12 (YoY)", 10.5, False, MUTED, False, 5, PP_ALIGN.LEFT, "Segoe UI"),
    ("· Rolling: 3-month and 6-month mean & std dev", 10.5, False, MUTED, False, 5, PP_ALIGN.LEFT, "Segoe UI"),
    ("· Momentum: MoM price diff & % change (lagged)", 10.5, False, MUTED, False, 5, PP_ALIGN.LEFT, "Segoe UI"),
    ("· Recursive roll: each predicted ŷ feeds t+1 as lag", 10.5, False, MUTED, False, 5, PP_ALIGN.LEFT, "Segoe UI"),
    ("", 6, False, INK, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("Best for: Tomatoes, Onions, all Pulses (except Gram)", 11, True, TEAL, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
], 6.95, 1.72, 5.75, 5.30)

# VS separator
rect(s, 6.55, 2.70, 0.025, 3.50, RULE)
textbox(s, "VS", 6.35, 3.70, 0.42, 0.42,
        size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER, font="Segoe UI")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7  ·  Results & The Potato Anomaly
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, SLATE)
hline(s, 0.55, 1.45, 12.23)
section_label(s, "Empirical Results")
slide_title(s, "Model Accuracy & The Potato Anomaly",
            "6-month holdout backtest. Accuracy % = 100 − MAPE.   RMSE & MAE measured in Rs/Kg.")

# Accuracy table — drawn manually as rows
headers = ["Commodity", "Category", "Prophet Acc %", "XGBoost Acc %", "Winner", "RMSE (Best)"]
col_widths = [2.50, 1.40, 1.65, 1.65, 1.40, 1.45]
col_lefts = [0.55]
for w in col_widths[:-1]:
    col_lefts.append(col_lefts[-1] + w + 0.04)

ROW_H = 0.42
TABLE_TOP = 1.60

# Header row
for i, (h, w, l) in enumerate(zip(headers, col_widths, col_lefts)):
    rect(s, l, TABLE_TOP, w, ROW_H, NAVY)
    textbox(s, h, l+0.08, TABLE_TOP+0.09, w-0.10, ROW_H-0.10,
            size=10, bold=True, color=WHITE, font="Segoe UI")

rows = [
    ("Pulse Moong",    "Pulses",     "94%",  "98%",  "XGBoost", "12.02",  TEAL,  False),
    ("Pulse Mash",     "Pulses",     "79%",  "94%",  "XGBoost", "29.83",  TEAL,  False),
    ("Pulse Masoor",   "Pulses",     "91%",  "93%",  "XGBoost", "17.75",  TEAL,  False),
    ("Pulse Gram",     "Pulses",     "90%",  "90%",  "Prophet", "32.02",  BLUE,  False),
    ("Tomatoes",       "Vegetables", "52%",  "80%",  "XGBoost", "59.52",  TEAL,  False),
    ("Onions",         "Vegetables", "42%",  "71%",  "XGBoost", "27.26",  TEAL,  False),
    ("Potatoes",       "Vegetables", "69%",  "−28%", "Prophet", "16.80",  BLUE,  True),
]

for ri, (comm, cat, prop_acc, xg_acc, winner, rmse, win_color, is_anomaly) in enumerate(rows):
    ty = TABLE_TOP + ROW_H + ri * ROW_H
    row_bg = RGBColor(255,251,235) if is_anomaly else (CARD_BG if ri % 2 == 0 else RGBColor(248,250,252))
    vals = [comm, cat, prop_acc, xg_acc, winner, rmse]
    for i, (val, w, l) in enumerate(zip(vals, col_widths, col_lefts)):
        rect(s, l, ty, w, ROW_H, row_bg, CARD_BDR, 0.5)
        cell_color = INK
        if i == 4:  cell_color = win_color
        if i == 3 and is_anomaly: cell_color = RGBColor(220,38,38)
        bold_cell = (i == 0 or i == 4)
        textbox(s, val, l+0.08, ty+0.09, w-0.10, ROW_H-0.10,
                size=10, bold=bold_cell, color=cell_color, font="Segoe UI")

# Anomaly callout box (right side, below table)
TABLE_BOT = TABLE_TOP + ROW_H * 8
rect(s, 0.55, TABLE_BOT + 0.10, 12.23, 1.42, RGBColor(255,247,237), RGBColor(253,186,116), 1.0, rounded=True)
multiline(s, [
    ("★  The Potato Anomaly  (XGBoost −28% Accuracy)", 11.5, True, AMBER, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("XGBoost's recursive multi-step roll compounded errors during a steep post-harvest crash (Rs 115 → Rs 31/kg), producing 127.9% MAPE. "
     "Prophet's Fourier sinusoidal wave naturally respected the seasonal trough — scoring 69%. "
     "This confirmed that dynamic model routing (no single winner) is mandatory.", 10.5, False, RGBColor(120,53,15), False, 4, PP_ALIGN.LEFT, "Segoe UI"),
], 0.80, TABLE_BOT + 0.18, 11.90, 1.25)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8  ·  Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, 0, 0, 13.333, 0.30, BLUE)

multiline(s, [
    ("Conclusion & Operational Impact", 24, True, WHITE, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
    ("ForecastBI transforms static government statistics into proactive, 12-month commodity price intelligence.", 12, False,
     RGBColor(147,197,253), False, 8, PP_ALIGN.LEFT, "Segoe UI"),
], 0.55, 0.42, 12.23, 1.10)

hline(s, 0.55, 1.48, 12.23, RGBColor(51,65,85))

# 4 impact tiles — 2x2 grid
tiles = [
    ("Automated Ingestion",
     "Eliminates manual collection. Hybrid XLSX/PDF parser compiles 60 months of PBS price bulletins automatically with full audit logging.",
     BLUE),
    ("Forward Visibility",
     "12-month forecasts per commodity give procurement teams and policymakers lead time to act before inflationary spikes occur.",
     TEAL),
    ("Empirical Model Routing",
     "Benchmark proves no single model wins universally. Dynamic routing (Prophet or XGBoost per item) maximizes forecast accuracy across all 7 commodities.",
     AMBER),
    ("Geospatial Price Intelligence",
     "City-level price hierarchy (Islamabad vs. Bannu) reveals arbitrage opportunities and regional logistics optimization potential across 17 commercial centers.",
     SKY),
]

positions = [(0.55, 1.60), (6.68, 1.60), (0.55, 4.38), (6.68, 4.38)]
W, H = 5.98, 2.58
for (tl, tt), (head, body, acc) in zip(positions, tiles):
    rect(s, tl, tt, W, H, RGBColor(18,32,60), RGBColor(51,65,85), 0.8, rounded=True)
    rect(s, tl, tt, W, 0.055, acc, rounded=False)
    multiline(s, [
        (head, 13, True, WHITE, False, 0, PP_ALIGN.LEFT, "Segoe UI"),
        (body, 10.5, False, RGBColor(203,213,225), False, 10, PP_ALIGN.LEFT, "Segoe UI"),
    ], tl+0.18, tt+0.13, W-0.28, H-0.20)

# Footer
rect(s, 0, 7.25, 13.333, 0.25, RGBColor(18,32,60))
textbox(s, "ForecastBI  ·  Pakistan Bureau of Statistics (PBS)  ·  Machine Learning & Power BI", 0.55, 7.27, 12.23, 0.22,
        size=8.5, bold=False, color=RGBColor(100,116,139), align=PP_ALIGN.CENTER, font="Segoe UI")

# ─── SAVE ───────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"[OK] Saved: {OUT}")
for dest in [REPORTS, ARTIFACT]:
    shutil.copy2(OUT, dest)
    print(f"[OK] Copied to: {dest}")
