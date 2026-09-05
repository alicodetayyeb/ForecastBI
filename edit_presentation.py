import sys
from pptx import Presentation

def replace_text(shape, new_text):
    if not shape.has_text_frame:
        return
    
    # capture font of first non-empty run
    font_name = None
    font_size = None
    font_bold = None
    font_italic = None
    font_color = None
    
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text.strip():
                font_name = r.font.name
                font_size = r.font.size
                font_bold = r.font.bold
                font_italic = r.font.italic
                try:
                    if r.font.color and r.font.color.type:
                        font_color = r.font.color.rgb
                except:
                    pass
                break
        if font_name is not None or font_size is not None:
            break

    # assign new text (this creates new paragraphs/runs)
    shape.text = new_text
    
    # re-apply formatting
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if font_name: r.font.name = font_name
            if font_size: r.font.size = font_size
            if font_bold is not None: r.font.bold = font_bold
            if font_italic is not None: r.font.italic = font_italic
            if font_color: r.font.color.rgb = font_color


prs = Presentation(r'D:\A- Commodity Price Coordinator\antigravity projects\ForecastBI\ForecastBI_Project_Presentation_sonnet_5.pptx')

replacements = {
    # Slide 1
    (0, 12): 'Pakistan Food Price Forecasting Dashboard',
    (0, 13): 'Data Collection • AI Forecasting • Interactive Dashboards',
    (0, 15): '5 Years of Data',
    (0, 17): '17 Major Cities',
    (0, 19): 'Smart AI Predictions',
    
    # Slide 2
    (1, 1): 'What is ForecastBI?',
    (1, 2): 'A smart tool to track and predict food prices in Pakistan to help manage food security.',
    (1, 7): 'Why It Matters',
    (1, 8): 'Food prices in Pakistan change rapidly. These sudden changes hurt everyday families the most. We need a way to see these changes coming.',
    (1, 12): 'What We Track',
    (1, 13): '3 Everyday Vegetables:\nTomatoes, Onions, Potatoes\n\n4 Essential Pulses:\nMoong, Masoor, Gram, Mash',
    (1, 17): 'Where We Track',
    (1, 18): 'We monitor prices across 17 major cities, including Karachi, Lahore, and Islamabad, using over 5 years of historical data.',
    
    # Slide 3
    (2, 1): 'The Challenge We Face',
    (2, 2): 'Why current methods are not enough to prevent price shocks.',
    (2, 7): 'Wild Price Swings',
    (2, 8): 'Vegetable prices can skyrocket unexpectedly. For example, tomato prices once jumped 174% in a single month. This makes planning very difficult.',
    (2, 12): 'Messy Data',
    (2, 13): 'Government price data is scattered across old Excel sheets and PDFs. It is hard to read, hard to combine, and full of missing pieces.',
    (2, 17): 'Acting Too Late',
    (2, 18): 'Because we cannot predict prices, the government and businesses can only react after a crisis happens, rather than preventing it.',
    
    # Slide 4
    (3, 0): 'SOLUTION',
    (3, 1): 'Our Solution: ForecastBI',
    (3, 2): 'An automated system that turns messy data into clear, future-looking predictions.',
    (3, 7): '1. Automated Data Collection',
    (3, 8): 'A software bot automatically downloads, reads, and cleans 5 years of government data every month with zero human effort.',
    (3, 12): '2. AI Forecasting Engine',
    (3, 13): 'We use two advanced AI models to predict future prices. The system automatically picks whichever model is most accurate for each specific food item.',
    (3, 17): '3. 12-Month Predictions',
    (3, 18): 'We generate realistic, 12-month future price forecasts to help businesses and the government plan ahead.',
    (3, 22): '4. Interactive Dashboards',
    (3, 23): 'We display everything on beautiful, easy-to-use dashboards so anyone can explore trends and compare city prices.',
    
    # Slide 5
    (4, 0): 'STRATEGY',
    (4, 1): 'How We Built It',
    (4, 2): 'A simple, 4-step process to ensure accuracy and reliability.',
    (4, 8): 'We carefully matched and cleaned data for different food items across 17 cities, fixing typos and missing values.',
    (4, 12): 'Finding Patterns',
    (4, 13): 'We taught the AI to look at historical trends, seasonal cycles, and recent price momentum.',
    (4, 17): 'Testing the Models',
    (4, 18): 'We trained the AI on older data and tested it by asking it to "predict" the last 6 months to see how well it did.',
    (4, 23): 'The system assigns the best-performing AI model to each specific food type to guarantee the best results.',
    
    # Slide 6
    (5, 1): 'Challenges We Overcame',
    (5, 2): 'Solving real-world data issues.',
    (5, 7): 'Hidden Website Data',
    (5, 10): 'The government website hid the download links inside website code, making them hard to find automatically.',
    (5, 12): 'We wrote a special script to read the hidden code and extract the links instantly.',
    (5, 19): 'The monthly reports came in three totally different file types: new Excel, old Excel, and PDF.',
    (5, 21): 'We built a smart reader that can read any of these formats line-by-line automatically.',
    (5, 28): 'Some months had broken web links or missing average prices.',
    (5, 30): 'Our system logs any broken links and automatically calculates missing averages using city data.',
    
    # Slide 7
    (6, 1): 'The Final Dashboards',
    (6, 2): 'Three interactive screens built for decision-makers.',
    (6, 8): '• Shows the latest prices, inflation rates, and 12-month forecasts.\n• Lets you click and filter by any food item (like Tomatoes or Onions).',
    (6, 13): '• A single chart showing both past prices (Blue) and future predictions (Orange).\n• Clearly shows when prices are expected to rise or fall.',
    (6, 18): '• Compares our AI models to show which one is winning.\n• Ranks the 17 cities from most expensive to least expensive.',
    
    # Slide 8
    (7, 0): 'RESULTS',
    (7, 1): 'Model Accuracy & The Potato Crash',
    (7, 2): 'Why we need more than one AI model to get things right.',
    (7, 5): 'Model Accuracy Scores',
    (7, 10): 'Key Insight: The Potato Anomaly',
    (7, 11): 'When potato prices crashed from Rs 115 to Rs 31, one of our AI models completely failed to understand it.\n\nWhy?\nThe first model just looked at recent price drops and thought the price would keep dropping forever.\nThe second model understood that potatoes follow a natural yearly harvest cycle, and correctly predicted the price would bounce back.\n\nTakeaway: We must use the right model for the right vegetable.',
    
    # Slide 9
    (8, 1): 'Conclusion & Impact',
    (8, 2): 'Moving from reacting to the past, to planning for the future.',
    (8, 7): 'For Decision Makers',
    (8, 8): 'Gives the government and distributors up to 12 months of warning before prices spike or crash.',
    (8, 12): 'For Businesses',
    (8, 13): 'Allows businesses to buy food when it is cheap and avoid buying during expensive peak seasons.',
    (8, 17): 'What is Next?',
    (8, 18): '• Adding weather and rainfall data to improve predictions.\n• Adding fuel and transport costs.\n• Tracking weekly prices instead of monthly.'
}

for i, slide in enumerate(prs.slides):
    for j, shape in enumerate(slide.shapes):
        key = (i, j)
        if key in replacements:
            replace_text(shape, replacements[key])

output_path = r'D:\A- Commodity Price Coordinator\antigravity projects\ForecastBI\ForecastBI_Project_Presentation_sonnet_5_Simple.pptx'
prs.save(output_path)
print(f"Saved successfully to {output_path}")

import shutil
shutil.copy2(output_path, r'C:\Users\sa\.gemini\antigravity-ide\brain\a572f10e-0c4b-4b65-a282-b8dfa98da6da\ForecastBI_Project_Presentation_Simple.pptx')
print("Artifact updated.")
