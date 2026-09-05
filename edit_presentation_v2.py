import sys
from pptx import Presentation
import shutil

def replace_text(shape, new_text):
    if not shape.has_text_frame:
        return
    
    font_name = None
    font_size = None
    font_bold = None
    font_italic = None
    font_color = None
    
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text.strip():
                font_name = r.font.name
                font_size = r.font.size
                font_bold = r.font.bold
                font_italic = r.font.italic
                try:
                    if r.font.color and r.font.color.type:
                        font_color = r.font.color.rgb
                except:
                    pass
                break
        if font_name is not None or font_size is not None:
            break

    shape.text = new_text
    
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if font_name: r.font.name = font_name
            if font_size: r.font.size = font_size
            if font_bold is not None: r.font.bold = font_bold
            if font_italic is not None: r.font.italic = font_italic
            if font_color: r.font.color.rgb = font_color


prs = Presentation(r'D:\A- Commodity Price Coordinator\antigravity projects\ForecastBI\ForecastBI_Project_Presentation_sonnet_5.pptx')

replacements = {
    # Slide 1 - Title
    (0, 11): 'ForecastBI: Data Science Project',
    (0, 12): 'End-to-End Time Series Forecasting & BI Analytics',
    (0, 13): 'Web Scraping • Machine Learning • Power BI',
    (0, 15): '5 Years of PBS Data',
    (0, 17): '17 Major Cities',
    (0, 19): 'Prophet vs. XGBoost',
    
    # Slide 2 - Intro
    (1, 0): 'INTRODUCTION',
    (1, 1): 'Project Background & Context',
    (1, 2): 'A data science project built to analyze and forecast government commodity prices.',
    (1, 7): 'The Goal',
    (1, 8): 'To build a complete data pipeline from scratch—extracting raw government data, training machine learning models, and visualizing the results.',
    (1, 12): 'Dataset Scope',
    (1, 13): '7 Commodities: Tomatoes, Onions, Potatoes, and 4 Pulses (Moong, Masoor, Gram, Mash).',
    (1, 17): 'Geographic Scope',
    (1, 18): 'Over 5 years of historical retail price data covering 17 major cities in Pakistan.',
    
    # Slide 3 - Problem Statement (Objectives)
    (2, 0): 'PROBLEM STATEMENT',
    (2, 1): 'Project Objectives & Scope',
    (2, 2): 'What we set out to build and the technical milestones we aimed to achieve.',
    (2, 7): '1. Data Acquisition',
    (2, 8): 'Programmatically scrape and compile 5 years of monthly historical data from the Pakistan Bureau of Statistics (PBS).',
    (2, 12): '2. Predictive Modeling',
    (2, 13): 'Engineer features and train machine learning models to forecast prices 12 months into the future.',
    (2, 17): '3. Data Visualization',
    (2, 18): 'Develop an interactive Power BI dashboard to showcase model accuracy, historical trends, and future forecasts.',
    
    # Slide 4 - Problems Faced (Dev challenges)
    (3, 0): 'PROBLEMS FACED',
    (3, 1): 'Data Preparation & Development Challenges',
    (3, 2): 'The real-world technical obstacles we encountered while building the data pipeline.',
    (3, 7): 'Challenge 1: Hidden Data Links',
    (3, 8): 'Standard web scrapers failed because the PBS website injected the file download links dynamically using client-side JavaScript.',
    (3, 12): 'Challenge 2: Chaotic File Formats',
    (3, 13): 'The historical reports were not standardized. We had to deal with modern Excel files (.xlsx), legacy Excel files (.xls), and unstructured PDFs.',
    (3, 17): 'Challenge 3: Missing & Broken Data',
    (3, 18): 'Several months had broken URLs, empty cells, or missing national average values that would break a standard ML pipeline.',
    (3, 22): '',
    (3, 23): '',
    
    # Slide 5 - Solution (How we overcame dev issues)
    (4, 0): 'SOLUTION',
    (4, 1): 'Overcoming the Development Issues',
    (4, 2): 'How we engineered our way out of the data collection roadblocks.',
    (4, 7): 'Bypassing JavaScript Execution',
    (4, 8): 'Instead of using slow headless browsers like Selenium, we used regular expressions and chompjs to parse the raw JavaScript variables directly in under 2 seconds.',
    (4, 12): 'Dual-Engine Document Parser',
    (4, 13): 'We built a custom parser combining openpyxl for Excel files and pypdf with regex tokenizers to extract data line-by-line from PDFs.',
    (4, 17): 'Automated Error Handling',
    (4, 18): 'We implemented detailed audit logging for missing files and wrote fallback logic to automatically calculate missing national averages from valid city data.',
    (4, 22): '',
    (4, 23): '',
    
    # Slide 6 - Strategy (ML/Tech approach)
    (5, 0): 'STRATEGY',
    (5, 1): 'Machine Learning & Technical Approach',
    (5, 2): 'The data science techniques used to process the data and generate forecasts.',
    (5, 7): 'Feature Engineering',
    (5, 9): 'Dates & Seasons',
    (5, 10): 'We transformed the raw dates into mathematical features (sine/cosine curves) so the models could understand yearly seasons.',
    (5, 11): 'Lags & Momentum',
    (5, 12): 'We added historical lag features (prices from 1, 3, and 6 months ago) and rolling averages.',
    (5, 16): 'Model Selection',
    (5, 18): 'Prophet vs. XGBoost',
    (5, 19): 'We chose Facebook Prophet (great for seasonal trends) and XGBoost (great for complex, non-linear relationships) to compete against each other.',
    (5, 20): '',
    (5, 21): '',
    (5, 25): 'Benchmarking Framework',
    (5, 27): '6-Month Holdout',
    (5, 28): 'We trained the models on the first 35 months and tested them on the final 6 months.',
    (5, 29): 'Evaluation Metric',
    (5, 30): 'We used MAPE (Mean Absolute Percentage Error) to determine the winner for each commodity.',
    
    # Slide 7 - Implementation
    (6, 0): 'IMPLEMENTATION',
    (6, 1): 'Bringing It All Together in Power BI',
    (6, 2): 'Connecting the Python backend to a front-end visualization layer.',
    (6, 7): 'Data Export & Star Schema',
    (6, 8): 'The Python pipeline exports the cleaned data and predictions as standardized CSV files. We structured these files in a relational format optimized for Power BI.',
    (6, 12): 'Dashboard Construction',
    (6, 13): 'We built a 3-page interactive report using a custom dark theme JSON file. The dashboard automatically reads the Python outputs to update its visuals.',
    (6, 17): 'Visual Analytics',
    (6, 18): 'The final dashboards feature interactive slicers, allowing users to explore historical trends, view future forecasts, and compare model accuracy dynamically.',
    
    # Slide 8 - Results
    (7, 0): 'RESULTS',
    (7, 1): 'Model Performance & Findings',
    (7, 2): 'The results of our 6-month machine learning backtest.',
    (7, 5): 'Model Accuracy Scores (100 - MAPE)',
    (7, 10): 'Key Insight: The Potato Anomaly',
    (7, 11): "XGBoost Dominated Pulses\nXGBoost won consistently on stable commodities like Pulses, achieving up to 98% accuracy because it effectively captured recent price momentum.\n\nThe Potato Anomaly\nWhen potato prices crashed from Rs 115 to Rs 31, XGBoost's predictions compounded errors and failed completely (-28% accuracy).\nFacebook Prophet succeeded here (69% accuracy) because its math is built on yearly sine waves, perfectly capturing the natural harvest cycle.\n\nConclusion: We proved that dynamic model routing (using different models for different items) was required for this project.",
    
    # Slide 9 - Conclusion
    (8, 0): 'CONCLUSION',
    (8, 1): 'Project Wrap-Up',
    (8, 2): 'Successfully demonstrating a full-stack data science workflow.',
    (8, 7): 'Pipeline Success',
    (8, 8): 'We successfully automated the extraction and cleaning of highly unstructured government data, replacing manual labor with a robust Python pipeline.',
    (8, 12): 'ML Validation',
    (8, 13): 'We proved the value of evaluating multiple algorithms. Relying on just one model would have caused massive errors in highly volatile vegetables.',
    (8, 17): 'Next Steps / Enhancements',
    (8, 18): '• Migrate the scraping pipeline to run automatically in the cloud.\n• Add weather and transportation cost data to the feature matrix.\n• Integrate the Power BI dashboard with a live database.'
}

for i, slide in enumerate(prs.slides):
    for j, shape in enumerate(slide.shapes):
        key = (i, j)
        if key in replacements:
            replace_text(shape, replacements[key])

output_path = r'D:\A- Commodity Price Coordinator\antigravity projects\ForecastBI\ForecastBI_Project_Presentation_sonnet_5_Tech.pptx'
prs.save(output_path)
print(f"Saved successfully to {output_path}")

shutil.copy2(output_path, r'C:\Users\sa\.gemini\antigravity-ide\brain\a572f10e-0c4b-4b65-a282-b8dfa98da6da\ForecastBI_Project_Presentation_Tech.pptx')
print("Artifact updated.")
